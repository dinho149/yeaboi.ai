"""Export a DeliveryReport as a native PowerPoint (.pptx) deck.

Mirrors the slide order of the HTML deck (presentation.py: title → executive
summary → metrics → one slide per outcome theme → highlights → delivered items →
thank-you) and colors every slide from the same palette the HTML deck uses
(themes.get_palette — built-ins or a custom palette from reporting_themes.json).

python-pptx is an *optional* dependency (the ``docs`` extra, shared with roadmap
intake): it is imported lazily inside the build function, and a missing install
degrades to returning None — the caller treats the .pptx as simply "not written".
All report text goes through python-pptx text frames, which write literal XML
text nodes — tracker-sourced strings are inert, no injection surface.
"""

from __future__ import annotations

import logging
import math
from pathlib import Path

from yeaboi.agent.state import DeliveryReport
from yeaboi.reporting.layout import (
    CARD_BULLET_PT,
    CARD_MIN_H_IN,
    CARD_TEXT_INSET_IN,
    CARD_TITLE_PT,
    LINE_SPACING,
    MARGIN_IN,
    SLIDE_H_IN,
    SLIDE_W_IN,
    TIGHT_CARDS_PER_SLIDE,
    card_height,
    est_lines,
    fit_bullets,
    plan_list_slides,
    plan_outcome_slides,
)
from yeaboi.reporting.style import FONT_PRESETS, FONT_SCALES, DeckStyle, resolve_color, summary_points
from yeaboi.reporting.themes import get_palette

logger = logging.getLogger(__name__)

# Geometry + text-fitting heuristics live in reporting/layout.py (shared with the
# planner, so estimates can never drift from the render); local aliases keep this
# module's historical names.
_SLIDE_W_IN = SLIDE_W_IN
_SLIDE_H_IN = SLIDE_H_IN
_MARGIN_IN = MARGIN_IN
_LINE_SPACING = LINE_SPACING
_est_lines = est_lines
_fit_bullets = fit_bullets

# Upper bound on delivered-item rows per table slide; the actual count is fitted
# to the slide from each row's estimated wrapped height (see est_lines).
_MAX_TABLE_ROWS = 15

# Longest delivered-item title shown in the table before an "…" cut.
_TITLE_CLIP = 110

# How many outcome-theme cards fit on one "compact" tight-fit slide (2×2 grid),
# mirroring the HTML deck's compact layout. Expand fit packs by content instead
# (see layout.plan_outcome_slides).
_CARDS_PER_SLIDE = TIGHT_CARDS_PER_SLIDE


def _emoji(report: DeliveryReport, slot: str, default: str = "") -> str:
    for s, e in report.emoji_theme:
        if s == slot and e:
            return e
    return default


def build_report_pptx(
    report: DeliveryReport, path: Path, *, theme: str = "midnight", style: DeckStyle | None = None
) -> Path | None:
    """Write ``report`` as a themed 16:9 .pptx deck at ``path``.

    ``style`` (see reporting/style.py) customizes fonts, colors, layout and optional
    sections; None means the neutral defaults — this function never reads the prefs
    file itself. Returns the written path, or None when python-pptx isn't installed
    (``uv sync --extra docs``) or the write fails — never raises.
    """
    try:
        from pptx import Presentation  # optional dependency: uv sync --extra docs
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches, Pt
    except ImportError:
        logger.info("python-pptx not installed — skipping .pptx export (install with: uv sync --extra docs)")
        return None

    style = style or DeckStyle()
    # Resolve "ask" for this non-interactive builder: it can never prompt, and
    # adding slides (never trimming content) is the safe default — the TUI export
    # flow resolves "ask" itself by offering the extra slides before calling in.
    fit_mode = style.content_fit if style.content_fit != "ask" else "expand"
    palette = get_palette(theme)
    colors = {role: RGBColor.from_string(value.lstrip("#")) for role, value in palette.items()}
    scale = FONT_SCALES.get(style.font_scale, 1.0)
    font_name = FONT_PRESETS.get(style.font_family, FONT_PRESETS["modern"])["pptx"]
    # Title/heading color overrides resolve to their historical role defaults.
    title_rgb = RGBColor.from_string(resolve_color(style.title_color, palette, palette["fg"]).lstrip("#"))
    heading_rgb = RGBColor.from_string(resolve_color(style.heading_color, palette, palette["accent"]).lstrip("#"))

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    def new_slide():
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg1"]
        return slide

    def add_text(slide, text, *, top, size, color, bold=False, left=_MARGIN_IN, width=None, height=1.0):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width or (_SLIDE_W_IN - 2 * left)), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(round(size * scale))
                run.font.bold = bold
                run.font.color.rgb = color
                run.font.name = font_name
        return box

    def add_bullets(slide, items, *, top, size=18, fit=True):
        # Fit to the box: the geometric budget matters as much as the user's cap —
        # six three-line bullets at "large" scale would run off the slide bottom.
        # fit=False renders items verbatim (the expand planner already paged them).
        box_w = _SLIDE_W_IN - 2 * _MARGIN_IN
        box_h = _SLIDE_H_IN - top - 0.6
        box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(top), Inches(box_w), Inches(box_h))
        tf = box.text_frame
        tf.word_wrap = True
        if fit:
            fitted = _fit_bullets(
                items, width_in=box_w, height_in=box_h, size_pt=size, scale=scale, max_items=style.max_bullets
            )
        else:
            fitted = [str(item) for item in items]
        for idx, item in enumerate(fitted):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            para.text = f"▸  {item}"
            para.space_after = Pt(10)
            for run in para.runs:
                run.font.size = Pt(round(size * scale))
                run.font.color.rgb = colors["fg"]
                run.font.name = font_name

    def add_card(slide, title, bullets, *, left, top, width, height, fitted=False):
        # Rounded-rect card mirroring the HTML deck's compact layout.
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["bg2"]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        # Shape text frames are middle-anchored by default, so overflow spills
        # both above AND below the card — anchor to the top and budget the
        # content so it stays inside instead.
        tf.vertical_anchor = MSO_ANCHOR.TOP
        # Insets + font sizes come from layout.py (card_height mirrors this math).
        tf.margin_left = Inches(CARD_TEXT_INSET_IN / 2)
        tf.margin_right = Inches(CARD_TEXT_INSET_IN / 2)
        tf.margin_top = Inches(0.1)  # 0.1 + 0.08 == layout.CARD_PAD_V_IN
        tf.margin_bottom = Inches(0.08)
        inner_w = width - CARD_TEXT_INSET_IN
        title_h = (
            _est_lines(title, inner_w, CARD_TITLE_PT, scale) * (_LINE_SPACING * CARD_TITLE_PT * scale / 72) + 6 / 72
        )
        body_h = height - 0.18 - title_h
        tf.paragraphs[0].text = title
        tf.paragraphs[0].space_after = Pt(6)
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(round(CARD_TITLE_PT * scale))
            run.font.bold = True
            run.font.color.rgb = heading_rgb
            run.font.name = font_name
        if fitted:  # the expand planner already paged these bullets to fit
            shown = [str(b) for b in bullets]
        else:
            shown = _fit_bullets(
                bullets,
                width_in=inner_w,
                height_in=body_h,
                size_pt=CARD_BULLET_PT,
                scale=scale,
                max_items=style.max_bullets,
                space_after_pt=4.0,
            )
        for item in shown:
            para = tf.add_paragraph()
            para.text = f"▸  {item}"
            para.space_after = Pt(4)
            for run in para.runs:
                run.font.size = Pt(round(CARD_BULLET_PT * scale))
                run.font.color.rgb = colors["fg"]
                run.font.name = font_name

    def add_duck(slide, *, left, top, height=0.5):
        """Stamp the yeaboi duck mark — best-effort, a missing asset is skipped."""
        from io import BytesIO

        from yeaboi.reporting.branding import duck_png

        data = duck_png()
        if data is None:
            return
        try:
            slide.shapes.add_picture(BytesIO(data), Inches(left), Inches(top), height=Inches(height))
        except Exception:  # noqa: BLE001 — branding is cosmetic, never kill the export
            logger.debug("reporting pptx: could not add branding image", exc_info=True)

    # -- Title slide --------------------------------------------------------
    slide = new_slide()
    dates = f"{report.period_start} to {report.period_end}".strip(" to")
    subtitle = report.period_label + (f"  ·  {dates}" if dates else "")
    if report.sprint_names:
        subtitle += f"  ·  {', '.join(report.sprint_names)}"
    add_text(slide, _emoji(report, "headline", "🚀"), top=1.2, size=48, color=colors["accent2"])
    add_text(slide, report.project_name or "Delivery Report", top=2.2, size=44, color=title_rgb, bold=True)
    add_text(slide, subtitle, top=3.4, size=18, color=colors["muted"])
    if report.headline:
        add_text(slide, report.headline, top=4.4, size=24, color=colors["accent2"], bold=True, height=2.0)
    add_duck(slide, left=_SLIDE_W_IN - 1.15, top=0.45)

    # -- Executive summary --------------------------------------------------
    if report.executive_summary:
        slide = new_slide()
        add_text(
            slide,
            f"{_emoji(report, 'summary', '📋')}  Executive summary",
            top=0.8,
            size=32,
            color=heading_rgb,
            bold=True,
        )
        # One prose paragraph is unreadable as a slide — split it into sentence
        # points, and step the size down until the estimate fits (never truncate).
        points = summary_points(report.executive_summary)
        if len(points) <= 1:
            add_text(slide, report.executive_summary, top=2.0, size=20, color=colors["fg"], height=4.5)
        else:
            body_w = _SLIDE_W_IN - 2 * _MARGIN_IN
            body_h = 4.7
            size = 18
            for candidate in (18, 16, 14):
                size = candidate
                est = sum(
                    _est_lines(p, body_w, candidate, scale) * (_LINE_SPACING * candidate * scale / 72) + 12 / 72
                    for p in points
                )
                if est <= body_h:
                    break
            box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(2.0), Inches(body_w), Inches(body_h))
            tf = box.text_frame
            tf.word_wrap = True
            for idx, point in enumerate(points):
                para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                para.text = f"▸  {point}"
                para.space_after = Pt(12)
                for run in para.runs:
                    run.font.size = Pt(round(size * scale))
                    run.font.color.rgb = colors["fg"]
                    run.font.name = font_name

    # -- Metrics ("By the numbers" tiles) -----------------------------------
    if report.metrics:
        slide = new_slide()
        add_text(
            slide,
            f"{_emoji(report, 'metrics', '📊')}  By the numbers",
            top=0.8,
            size=32,
            color=heading_rgb,
            bold=True,
        )
        tile_w = 2.8
        per_row = max(1, int((_SLIDE_W_IN - 2 * _MARGIN_IN) // tile_w))
        for idx, (label, value) in enumerate(report.metrics):
            row, col = divmod(idx, per_row)
            left = _MARGIN_IN + col * tile_w
            top = 2.2 + row * 1.8
            if top > _SLIDE_H_IN - 1.5:
                break
            add_text(slide, str(value), top=top, size=40, color=colors["accent2"], bold=True, left=left, width=tile_w)
            add_text(slide, str(label), top=top + 0.9, size=14, color=colors["muted"], left=left, width=tile_w)
        if report.supporting_signals and style.include_signals:
            from yeaboi.reporting.context import signals_sentence

            sentence = signals_sentence(report.supporting_signals)
            if sentence:
                # Same corroboration footnote as the HTML deck's metrics slide.
                add_text(slide, sentence, top=_SLIDE_H_IN - 1.0, size=12, color=colors["muted"])

    # -- Outcome themes: one slide each, or grouped as cards (compact layout) --
    if style.layout == "compact" and report.themes:
        gap = 0.3
        card_w = (_SLIDE_W_IN - 2 * _MARGIN_IN - gap) / 2
        content_h = _SLIDE_H_IN - 2.0 - 0.5  # heading band above, breathing room below
        if fit_mode == "tight":
            # Fixed 2×2 grid; card content is trimmed to the grid geometry.
            chunks = [report.themes[i : i + _CARDS_PER_SLIDE] for i in range(0, len(report.themes), _CARDS_PER_SLIDE)]
            for chunk_idx, chunk in enumerate(chunks, start=1):
                slide = new_slide()
                heading = "Outcomes" if len(chunks) == 1 else f"Outcomes ({chunk_idx}/{len(chunks)})"
                add_text(
                    slide,
                    f"{_emoji(report, 'themes', '🧩')}  {heading}",
                    top=0.8,
                    size=32,
                    color=heading_rgb,
                    bold=True,
                )
                # A 1–2 card chunk gets full-height cards instead of a lone quarter-size one.
                rows = math.ceil(len(chunk) / 2)
                card_h = (content_h - (rows - 1) * gap) / rows
                for idx, (ttitle, outcomes) in enumerate(chunk):
                    row, col = divmod(idx, 2)
                    add_card(
                        slide,
                        ttitle,
                        outcomes,
                        left=_MARGIN_IN + col * (card_w + gap),
                        top=2.0 + row * (card_h + gap),
                        width=card_w,
                        height=card_h,
                    )
        else:
            # Expand fit: the planner keeps every bullet (long themes paginate
            # into "(cont.)" cards) and packs content-sized cards into slides;
            # heights come from the same layout.card_height the planner used.
            plan = plan_outcome_slides(report.themes, scale=scale, max_bullets=style.max_bullets)
            for slide_idx, slide_plan in enumerate(plan, start=1):
                slide = new_slide()
                heading = "Outcomes" if len(plan) == 1 else f"Outcomes ({slide_idx}/{len(plan)})"
                add_text(
                    slide,
                    f"{_emoji(report, 'themes', '🧩')}  {heading}",
                    top=0.8,
                    size=32,
                    color=heading_rgb,
                    bold=True,
                )
                cards = slide_plan.cards
                if len(cards) == 1 and cards[0].full_width:
                    # A lone card spans the full width at its content height
                    # instead of leaving an empty column.
                    card = cards[0]
                    full_w = _SLIDE_W_IN - 2 * _MARGIN_IN
                    h = max(CARD_MIN_H_IN, card_height(card.title, card.bullets, width_in=full_w, scale=scale))
                    add_card(
                        slide,
                        card.title,
                        card.bullets,
                        left=_MARGIN_IN,
                        top=2.0,
                        width=full_w,
                        height=min(h, content_h),
                        fitted=True,
                    )
                    continue
                y = 2.0
                for row_start in range(0, len(cards), 2):
                    row = cards[row_start : row_start + 2]
                    heights = [
                        max(CARD_MIN_H_IN, card_height(c.title, c.bullets, width_in=card_w, scale=scale)) for c in row
                    ]
                    for col, (card, h) in enumerate(zip(row, heights)):
                        add_card(
                            slide,
                            card.title,
                            card.bullets,
                            left=_MARGIN_IN + col * (card_w + gap),
                            top=y,
                            width=card_w,
                            height=min(h, content_h),
                            fitted=True,
                        )
                    y += max(heights) + gap
    else:
        for ttitle, outcomes in report.themes:
            if fit_mode == "tight":
                pages = [(ttitle, tuple(outcomes))]  # single slide; add_bullets trims
            else:
                pages = plan_list_slides(ttitle, outcomes, scale=scale, max_bullets=style.max_bullets)
            for page_title, page_items in pages:
                slide = new_slide()
                add_text(
                    slide,
                    f"{_emoji(report, 'themes', '🧩')}  {page_title}",
                    top=0.8,
                    size=32,
                    color=heading_rgb,
                    bold=True,
                )
                add_bullets(slide, page_items, top=2.0, fit=fit_mode == "tight")

    # -- Production ----------------------------------------------------------
    # Its own slide rather than a second footnote on "By the numbers": that line
    # already claims corroboration, and an incident is not corroboration.
    if getattr(report, "ops_signals", ()) and style.include_production:
        from yeaboi.ops.signals import describe
        from yeaboi.reporting.context import OPS_EMOJI, ops_sentence

        slide = new_slide()
        add_text(slide, f"{OPS_EMOJI}  Production", top=0.8, size=32, color=heading_rgb, bold=True)
        add_bullets(
            slide,
            [describe(sig) for sig in report.ops_signals][: style.max_bullets],
            top=2.0,
            fit=fit_mode == "tight",
        )
        sentence = ops_sentence(report.ops_signals)
        if sentence:
            add_text(
                slide,
                f"{sentence}. Team-wide, and not attributed to anyone.",
                top=_SLIDE_H_IN - 1.0,
                size=12,
                color=colors["muted"],
            )

    # -- Highlights ----------------------------------------------------------
    if report.highlights and style.include_highlights:
        if fit_mode == "tight":
            pages = [("Highlights", tuple(report.highlights))]
        else:
            pages = plan_list_slides("Highlights", report.highlights, scale=scale, max_bullets=style.max_bullets)
        for page_title, page_items in pages:
            slide = new_slide()
            add_text(
                slide,
                f"{_emoji(report, 'highlights', '⭐')}  {page_title}",
                top=0.8,
                size=32,
                color=heading_rgb,
                bold=True,
            )
            add_bullets(slide, page_items, top=2.0, fit=fit_mode == "tight")

    # -- Delivered items table -----------------------------------------------
    if report.delivered_items and style.include_items_table:
        slide = new_slide()
        add_text(slide, "Delivered items", top=0.8, size=32, color=heading_rgb, bold=True)
        items = list(report.delivered_items)

        def _short_title(title: str) -> str:
            return title if len(title) <= _TITLE_CLIP else title[: _TITLE_CLIP - 1].rstrip() + "…"

        # Fit rows to the slide: pptx row heights are minimums — wrapped titles
        # grow them — so accumulate estimated heights instead of a fixed count.
        col_widths = (1.5, 6.5, 1.3, 2.23)  # Key / Title / Status / Assignee (sums to the usable width)
        line_h = _LINE_SPACING * 11 * scale / 72
        table_top = 1.9
        budget = 4.6  # inches available under the heading, before the "more" line
        used = 0.35  # header row
        shown: list = []
        for item in items[:_MAX_TABLE_ROWS]:
            est = _est_lines(_short_title(item.title), 6.3, 11, scale) * line_h + 0.12
            reserve = 0.3 if len(shown) + 1 < len(items) else 0.0  # keep room for the "more" line
            if shown and used + est + reserve > budget:
                break
            used += est
            shown.append(item)
        table_shape = slide.shapes.add_table(
            len(shown) + 1,
            4,
            Inches(_MARGIN_IN),
            Inches(table_top),
            Inches(_SLIDE_W_IN - 2 * _MARGIN_IN),
            Inches(used),
        )
        # Swap the default banded-blue table style for "No Style, No Grid" so the
        # explicit dark-theme fills below are all that shows. Best-effort — this
        # pokes python-pptx XML internals, and cosmetics must never kill the export.
        try:
            tbl_props = table_shape._element.graphic.graphicData.tbl[0]
            tbl_props[-1].text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
        except Exception:  # noqa: BLE001 — table keeps the default style, still readable
            logger.debug("reporting pptx: could not neutralise table style", exc_info=True)
        table = table_shape.table
        for col, width in enumerate(col_widths):
            table.columns[col].width = Inches(width)
        for col, heading in enumerate(("Key", "Title", "Status", "Assignee")):
            table.cell(0, col).text = heading
        for row, item in enumerate(shown, start=1):
            table.cell(row, 0).text = item.key
            table.cell(row, 1).text = _short_title(item.title)
            table.cell(row, 2).text = item.status
            table.cell(row, 3).text = item.assignee
        for row_idx, table_row in enumerate(table.rows):
            for cell in table_row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["bg2"]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(round(11 * scale))
                        run.font.name = font_name
                        run.font.bold = row_idx == 0
                        run.font.color.rgb = heading_rgb if row_idx == 0 else colors["fg"]
        if len(items) > len(shown):
            add_text(
                slide,
                f"… and {len(items) - len(shown)} more",
                top=table_top + used + 0.1,
                size=14,
                color=colors["muted"],
            )

    # -- Thank-you ------------------------------------------------------------
    if style.include_thanks:
        slide = new_slide()
        add_text(slide, _emoji(report, "thanks", "🙌"), top=2.2, size=54, color=colors["accent2"])
        add_text(slide, "Thank you", top=3.5, size=40, color=title_rgb, bold=True)
        if report.project_name:
            add_text(slide, report.project_name, top=4.6, size=18, color=colors["muted"])
        # Small duck + wordmark, centered under the message (clear of the
        # bottom-corner footer/slide-number chrome).
        add_duck(slide, left=_SLIDE_W_IN / 2 - 0.25, top=5.7)
        add_text(
            slide,
            "made with yeaboi.ai",
            top=6.3,
            size=10,
            color=colors["muted"],
            left=_SLIDE_W_IN / 2 - 0.65,
            width=1.6,
            height=0.3,
        )

    # -- Per-slide chrome: custom footer + slide numbers ----------------------
    if style.footer_text or style.slide_numbers:
        for idx, slide in enumerate(prs.slides, start=1):
            if style.footer_text:
                add_text(
                    slide,
                    style.footer_text,
                    top=_SLIDE_H_IN - 0.45,
                    size=10,
                    color=colors["muted"],
                    width=8.0,
                    height=0.35,
                )
            if style.slide_numbers:
                add_text(
                    slide,
                    str(idx),
                    top=_SLIDE_H_IN - 0.45,
                    size=10,
                    color=colors["muted"],
                    left=_SLIDE_W_IN - 1.2,
                    width=0.6,
                    height=0.35,
                )

    try:
        prs.save(str(path))
    except OSError as e:
        logger.warning("reporting pptx export failed: %s", e)
        return None
    logger.info("Reporting .pptx exported: %s (theme=%s, %d slides)", path, theme, len(prs.slides))
    return path
