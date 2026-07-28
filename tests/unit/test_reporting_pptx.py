"""Unit tests for the native PowerPoint (.pptx) reporting exporter."""

import builtins

import pytest

from yeaboi.agent.state import DeliveredItem, DeliveryReport
from yeaboi.reporting.pptx_export import build_report_pptx


def _report(n_items: int = 3) -> DeliveryReport:
    items = tuple(
        DeliveredItem(key=f"P-{i}", title=f"feature {i}", status="Done", source="jira", assignee="Ada")
        for i in range(n_items)
    )
    return DeliveryReport(
        period_label="Last sprint",
        period_start="2026-07-01",
        period_end="2026-07-14",
        project_name="Demo",
        headline="A strong sprint.",
        executive_summary="We shipped a lot of things.",
        themes=(("Security", ("SSO shipped", "MFA rolled out")),),
        highlights=("SSO live",),
        metrics=(("Items delivered", str(n_items)), ("Contributors", "1")),
        delivered_items=items,
        emoji_theme=(("headline", "🚀"),),
        generated_at="2026-07-14",
    )


class TestMissingDependency:
    def test_returns_none_when_pptx_not_installed(self, tmp_path, monkeypatch):
        real_import = builtins.__import__

        def _no_pptx(name, *args, **kwargs):
            if name == "pptx" or name.startswith("pptx."):
                raise ImportError("No module named 'pptx'")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _no_pptx)
        assert build_report_pptx(_report(), tmp_path / "deck.pptx", theme="midnight") is None
        assert not (tmp_path / "deck.pptx").exists()


class TestDeckStructure:
    """Real python-pptx assertions — skipped when the docs extra isn't installed."""

    @pytest.fixture(autouse=True)
    def _pptx(self):
        return pytest.importorskip("pptx", reason="docs extra not installed")

    def test_writes_deck_with_expected_slides(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", theme="aurora")
        assert path is not None and path.exists()

        from pptx import Presentation

        prs = Presentation(str(path))
        # Title, summary, metrics, 1 theme, highlights, delivered items, thank-you.
        assert len(prs.slides) == 7
        texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ]
        joined = "\n".join(texts)
        assert "Demo" in joined
        assert "A strong sprint." in joined
        assert "Executive summary" in joined
        assert "Security" in joined
        assert "SSO shipped" in joined
        assert "Thank you" in joined

    def test_sixteen_nine_geometry_and_theme_background(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", theme="sunset")
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(str(path))
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)
        # Every slide gets the palette's bg1 as a solid background fill.
        from yeaboi.reporting.themes import get_palette

        expected = get_palette("sunset")["bg1"].lstrip("#").upper()
        for slide in prs.slides:
            assert str(slide.background.fill.fore_color.rgb) == expected

    def test_items_table_caps_rows_with_overflow_note(self, tmp_path):
        path = build_report_pptx(_report(n_items=40), tmp_path / "deck.pptx")
        from pptx import Presentation

        prs = Presentation(str(path))
        tables = [shape.table for slide in prs.slides for shape in slide.shapes if shape.has_table]
        assert len(tables) == 1
        # Rows are height-fitted to the slide (upper-bounded by _MAX_TABLE_ROWS + header).
        n_rows = len(tables[0].rows)
        assert 2 <= n_rows <= 16
        texts = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert f"and {40 - (n_rows - 1)} more" in texts

    def test_empty_report_still_builds(self, tmp_path):
        path = build_report_pptx(DeliveryReport(period_label="Last week"), tmp_path / "deck.pptx")
        assert path is not None and path.exists()

    def test_unknown_theme_falls_back_to_midnight(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", theme="does-not-exist")
        from pptx import Presentation

        from yeaboi.reporting.themes import BUILTIN_PALETTES

        prs = Presentation(str(path))
        expected = BUILTIN_PALETTES["midnight"]["bg1"].lstrip("#").upper()
        assert str(prs.slides[0].background.fill.fore_color.rgb) == expected


class TestSupportingSignalsFootnote:
    @pytest.fixture(autouse=True)
    def _pptx(self):
        return pytest.importorskip("pptx", reason="docs extra not installed")

    def _signals_report(self):
        from dataclasses import replace

        from yeaboi.agent.state import SupportingSignal

        return replace(
            _report(),
            supporting_signals=(
                SupportingSignal(kind="pull_requests", source="github", count=24),
                SupportingSignal(kind="doc_updates", source="notion", count=5),
            ),
        )

    def _all_text(self, path):
        from pptx import Presentation

        prs = Presentation(str(path))
        return "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )

    def test_metrics_slide_footnote_present(self, tmp_path):
        path = build_report_pptx(self._signals_report(), tmp_path / "deck.pptx")
        assert "Corroborated by 24 merged PRs and 5 doc updates" in self._all_text(path)

    def test_no_signals_no_footnote(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx")
        assert "Corroborated by" not in self._all_text(path)

    def test_include_signals_false_drops_footnote(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(self._signals_report(), tmp_path / "deck.pptx", style=DeckStyle(include_signals=False))
        assert "Corroborated by" not in self._all_text(path)


class TestDeckStyle:
    """DeckStyle customization — fonts, colors, layout, toggles, chrome."""

    @pytest.fixture(autouse=True)
    def _pptx(self):
        return pytest.importorskip("pptx", reason="docs extra not installed")

    def _runs(self, path):
        from pptx import Presentation

        prs = Presentation(str(path))
        return [
            run
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]

    def _all_text(self, path):
        return "\n".join(run.text for run in self._runs(path))

    def _slides(self, path):
        from pptx import Presentation

        return Presentation(str(path)).slides

    def test_default_style_is_identity(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        a = build_report_pptx(_report(), tmp_path / "a.pptx")
        b = build_report_pptx(_report(), tmp_path / "b.pptx", style=DeckStyle())
        assert len(self._slides(a)) == len(self._slides(b)) == 7
        assert self._all_text(a) == self._all_text(b)

    def test_font_preset_sets_typeface_on_every_run(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(_report(), tmp_path / "deck.pptx", style=DeckStyle(font_family="classic"))
        names = {run.font.name for run in self._runs(path)}
        assert names == {"Georgia"}

    def test_font_scale_multiplies_sizes(self, tmp_path):
        from pptx.util import Pt

        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(_report(), tmp_path / "deck.pptx", style=DeckStyle(font_scale="large"))
        title_run = next(run for run in self._runs(path) if run.text == "Demo")
        assert title_run.font.size == Pt(round(44 * 1.15))

    def test_title_and_heading_color_overrides(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle
        from yeaboi.reporting.themes import get_palette

        style = DeckStyle(title_color="#ff0000", heading_color="accent2")
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", theme="midnight", style=style)
        runs = self._runs(path)
        title_run = next(run for run in runs if run.text == "Demo")
        assert str(title_run.font.color.rgb) == "FF0000"
        heading_run = next(run for run in runs if "Executive summary" in run.text)
        assert str(heading_run.font.color.rgb) == get_palette("midnight")["accent2"].lstrip("#").upper()

    def test_toggles_drop_optional_slides(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        style = DeckStyle(
            include_items_table=False, include_signals=False, include_highlights=False, include_thanks=False
        )
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", style=style)
        # Title, summary, metrics, 1 theme — appendix/highlights/thanks all gone.
        assert len(self._slides(path)) == 4
        text = self._all_text(path)
        assert "Thank you" not in text and "Highlights" not in text and "Delivered items" not in text

    def test_compact_layout_groups_themes_as_cards(self, tmp_path):
        from dataclasses import replace

        from yeaboi.reporting.style import DeckStyle

        report = replace(_report(), themes=tuple((f"Theme {i}", (f"outcome {i}",)) for i in range(5)))
        detailed = build_report_pptx(report, tmp_path / "a.pptx")
        compact = build_report_pptx(report, tmp_path / "b.pptx", style=DeckStyle(layout="compact", content_fit="tight"))
        # 5 per-theme slides collapse into 2 card slides (4 cards + 1 card).
        assert len(self._slides(detailed)) - len(self._slides(compact)) == 3
        text = self._all_text(compact)
        assert "Outcomes (1/2)" in text and "Outcomes (2/2)" in text
        assert "Theme 0" in text and "Theme 4" in text and "outcome 2" in text

    def test_compact_expand_coalesces_short_themes(self, tmp_path):
        from dataclasses import replace

        from yeaboi.reporting.style import DeckStyle

        report = replace(_report(), themes=tuple((f"Theme {i}", (f"outcome {i}",)) for i in range(5)))
        path = build_report_pptx(report, tmp_path / "deck.pptx", style=DeckStyle(layout="compact"))
        text = self._all_text(path)
        # Content-sized cards let all 5 short themes share one plain-titled slide.
        assert "Outcomes" in text and "Outcomes (1/" not in text
        for i in range(5):
            assert f"outcome {i}" in text
        assert "… and" not in text

    def test_max_bullets_caps_with_overflow_marker(self, tmp_path):
        from dataclasses import replace

        from yeaboi.reporting.style import DeckStyle

        report = replace(_report(), themes=(("Big", tuple(f"item {i}" for i in range(8))),))
        path = build_report_pptx(report, tmp_path / "deck.pptx", style=DeckStyle(max_bullets=3, content_fit="tight"))
        text = self._all_text(path)
        assert "item 2" in text and "item 3" not in text
        assert "… and 5 more" in text

    def test_expand_paginates_detailed_list_instead_of_trimming(self, tmp_path):
        from dataclasses import replace

        from yeaboi.reporting.style import DeckStyle

        report = replace(_report(), themes=(("Big", tuple(f"item {i}" for i in range(8))),))
        path = build_report_pptx(report, tmp_path / "deck.pptx", style=DeckStyle(max_bullets=3))
        text = self._all_text(path)
        # max_bullets acts as a page size: 8 items over 3 slides, nothing dropped.
        assert "Big (1/3)" in text and "Big (3/3)" in text
        for i in range(8):
            assert f"item {i}" in text
        assert "… and" not in text

    def test_footer_and_slide_numbers_on_every_slide(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        style = DeckStyle(footer_text="ACME Corp — internal", slide_numbers=True)
        path = build_report_pptx(_report(), tmp_path / "deck.pptx", style=style)
        slides = self._slides(path)
        for idx, slide in enumerate(slides, start=1):
            texts = [shape.text_frame.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
            assert "ACME Corp — internal" in texts
            assert str(idx) in texts


class TestBranding:
    """The yeaboi duck mark on the title + thank-you slides."""

    @pytest.fixture(autouse=True)
    def _pptx(self):
        return pytest.importorskip("pptx", reason="docs extra not installed")

    def _pictures_per_slide(self, path):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return [
            [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            for slide in Presentation(str(path)).slides
        ]

    def test_duck_on_title_and_thanks_slides(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx")
        pics = self._pictures_per_slide(path)
        assert len(pics[0]) == 1  # title slide
        assert len(pics[-1]) == 1  # thank-you slide
        assert all(not p for p in pics[1:-1])  # nowhere else — the mark stays subtle

    def test_thanks_slide_carries_wordmark(self, tmp_path):
        from pptx import Presentation

        path = build_report_pptx(_report(), tmp_path / "deck.pptx")
        last = list(Presentation(str(path)).slides)[-1]
        texts = [sh.text_frame.text for sh in last.shapes if getattr(sh, "has_text_frame", False)]
        assert "made with yeaboi.ai" in texts

    def test_missing_duck_asset_still_exports(self, tmp_path, monkeypatch):
        monkeypatch.setattr("yeaboi.reporting.branding.duck_png", lambda: None)
        path = build_report_pptx(_report(), tmp_path / "deck.pptx")
        assert path is not None
        assert all(not pics for pics in self._pictures_per_slide(path))


_LONG_OUTCOME = (
    "Delivered a substantial improvement to the platform that touches several critical "
    "systems, hardens the security posture across every region, and reduces operational risk."
)


class TestVisualFit:
    """Overflow fixes — card fitting, table fitting, structured summary."""

    @pytest.fixture(autouse=True)
    def _pptx(self):
        return pytest.importorskip("pptx", reason="docs extra not installed")

    def _dense_report(self, n_themes: int = 4) -> DeliveryReport:
        from dataclasses import replace

        themes = tuple(
            (f"Long Strategic Theme Title Number {i} With Extra Words", tuple(_LONG_OUTCOME for _ in range(6)))
            for i in range(n_themes)
        )
        return replace(_report(), themes=themes)

    def _cards(self, path):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return [
            shape
            for slide in Presentation(str(path)).slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]

    def test_card_text_top_anchored_and_fitted(self, tmp_path):
        from pptx.enum.text import MSO_ANCHOR

        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(
            self._dense_report(), tmp_path / "deck.pptx", style=DeckStyle(layout="compact", content_fit="tight")
        )
        cards = self._cards(path)
        assert cards, "compact layout must render rounded-rect cards"
        for card in cards:
            tf = card.text_frame
            assert tf.vertical_anchor == MSO_ANCHOR.TOP
            bullets = [p for p in tf.paragraphs if p.text.startswith("▸")]
            # Six long outcomes cannot fit a quarter-slide card — the budget trims
            # them and the trailing marker accounts for the rest.
            assert len(bullets) < 6
            assert "more" in bullets[-1].text

    def test_expand_cards_keep_every_bullet_without_markers(self, tmp_path):
        from dataclasses import replace

        from pptx.enum.text import MSO_ANCHOR

        from yeaboi.reporting.style import DeckStyle

        # 8 outcomes per theme exceed the max_bullets=6 page size, forcing "(cont.)".
        themes = tuple((f"Big Theme {i}", tuple(_LONG_OUTCOME for _ in range(8))) for i in range(2))
        report = replace(_report(), themes=themes)
        path = build_report_pptx(report, tmp_path / "deck.pptx", style=DeckStyle(layout="compact"))
        cards = self._cards(path)
        assert cards
        bullets = [p.text for c in cards for p in c.text_frame.paragraphs if p.text.startswith("▸")]
        # 2 themes × 8 long outcomes all render, spread across "(cont.)" cards.
        assert len(bullets) == 16
        assert all("more" not in b for b in bullets)
        assert any("(cont.)" in c.text_frame.paragraphs[0].text for c in cards)
        assert all(c.text_frame.vertical_anchor == MSO_ANCHOR.TOP for c in cards)

    def test_small_final_chunk_gets_taller_cards(self, tmp_path):
        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(
            self._dense_report(n_themes=5),
            tmp_path / "deck.pptx",
            style=DeckStyle(layout="compact", content_fit="tight"),
        )
        heights = sorted({card.height for card in self._cards(path)})
        assert len(heights) == 2  # 2×2 grid cards + the lone full-height card
        assert heights[-1] > heights[0] * 1.8

    def test_expand_lone_card_slides_render_full_width(self, tmp_path):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches

        from yeaboi.reporting.style import DeckStyle

        path = build_report_pptx(
            self._dense_report(n_themes=5), tmp_path / "deck.pptx", style=DeckStyle(layout="compact")
        )
        slides = Presentation(str(path)).slides
        per_slide = [[sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE] for slide in slides]
        lone = [cards[0] for cards in per_slide if len(cards) == 1]
        assert lone, "dense 5-theme expand plan should leave at least one lone-card slide"
        for card in lone:
            assert card.width > Inches(8)  # full-width, not the half-slide card
        # Cards are content-sized: no card overruns the content band.
        for cards in per_slide:
            for card in cards:
                assert card.top + card.height <= Inches(7.5 - 0.5) + 9525  # 0.001" tolerance

    def test_table_column_widths_and_dark_fills(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches

        from yeaboi.reporting.themes import get_palette

        path = build_report_pptx(_report(n_items=5), tmp_path / "deck.pptx", theme="midnight")
        prs = Presentation(str(path))
        tables = [shape.table for slide in prs.slides for shape in slide.shapes if shape.has_table]
        assert len(tables) == 1
        table = tables[0]
        # Title column dominates; Key no longer wastes a quarter of the slide.
        assert table.columns[1].width == Inches(6.5)
        assert table.columns[1].width > 3 * table.columns[0].width
        bg2 = get_palette("midnight")["bg2"].lstrip("#").upper()
        for row_idx, row in enumerate(table.rows):
            for cell in row.cells:
                assert str(cell.fill.fore_color.rgb) == bg2
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.bold == (row_idx == 0)

    def test_long_titles_shrink_row_count(self, tmp_path):
        from dataclasses import replace

        from pptx import Presentation

        long_title = "An extremely long delivered item title " * 4
        items = tuple(
            DeliveredItem(key=f"P-{i}", title=long_title, status="Done", source="jira", assignee="Ada")
            for i in range(40)
        )
        path = build_report_pptx(replace(_report(), delivered_items=items), tmp_path / "deck.pptx")
        prs = Presentation(str(path))
        tables = [shape.table for slide in prs.slides for shape in slide.shapes if shape.has_table]
        n_rows = len(tables[0].rows)
        assert n_rows < 16  # wrapped titles eat the height budget → fewer rows fit
        texts = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert f"and {40 - (n_rows - 1)} more" in texts
        # Cell titles are clipped so a single row can't wrap indefinitely.
        assert all(len(cell.text) <= 110 for row in tables[0].rows for cell in row.cells)

    def test_multi_sentence_summary_renders_points(self, tmp_path):
        from dataclasses import replace

        summary = (
            "Over the past month the team focused on locking down access across the estate. "
            "Monitoring coverage expanded to every cloud account and identity system we operate. "
            "A large housekeeping programme removed legacy accounts and tightened permissions. "
            "All of it was backed by one hundred merged pull requests."
        )
        path = build_report_pptx(replace(_report(), executive_summary=summary), tmp_path / "deck.pptx")
        from pptx import Presentation

        prs = Presentation(str(path))
        summary_slide = list(prs.slides)[1]
        points = [
            para.text
            for shape in summary_slide.shapes
            if getattr(shape, "has_text_frame", False)
            for para in shape.text_frame.paragraphs
            if para.text.startswith("▸")
        ]
        assert len(points) == 4
        assert points[0].endswith("estate.")

    def test_single_sentence_summary_keeps_plain_text(self, tmp_path):
        path = build_report_pptx(_report(), tmp_path / "deck.pptx")
        from pptx import Presentation

        prs = Presentation(str(path))
        summary_slide = list(prs.slides)[1]
        texts = [shape.text_frame.text for shape in summary_slide.shapes if getattr(shape, "has_text_frame", False)]
        assert "We shipped a lot of things." in texts  # no bullet glyph added
